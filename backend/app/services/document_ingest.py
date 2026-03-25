"""Document ingestion pipeline — PDF, PPTX, video, text/markdown.

Processes documents into chunks, embeds them, and stores in Pinecone.
"""

import io
import logging
import os
import tempfile
import uuid
from datetime import datetime, timezone
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation

from app.config import settings
from app.models.kb_source import KBSource
from app.utils.chunking import (
    TextChunk,
    ImageChunk,
    chunk_by_sections,
    chunk_markdown,
    chunk_plain_text,
)
from app.utils.embeddings import embed_text, embed_image, image_to_base64

logger = logging.getLogger(__name__)


def _ensure_dir(path: str) -> str:
    os.makedirs(path, exist_ok=True)
    return path


def _save_image(image_bytes: bytes, base_dir: str, name: str) -> str:
    """Save image bytes to disk, return the file path."""
    _ensure_dir(base_dir)
    path = os.path.join(base_dir, name)
    with open(path, "wb") as f:
        f.write(image_bytes)
    return path


# ---------------------------------------------------------------------------
# PDF Ingestion
# ---------------------------------------------------------------------------

async def ingest_pdf(
    file_path: str,
    source_id: str,
    pinecone_index,
) -> int:
    """Ingest a PDF: extract text chunks + images, embed, store in Pinecone.

    Returns the total number of chunks created.
    """
    doc = fitz.open(file_path)
    all_vectors = []
    chunk_count = 0
    image_dir = os.path.join(settings.UPLOAD_DIR, "images", source_id)

    for page_num in range(len(doc)):
        page = doc[page_num]

        # --- Text extraction ---
        text = page.get_text("text")
        if text.strip():
            text_chunks = chunk_by_sections(text, page_num=page_num + 1)
            for tc in text_chunks:
                chunk_id = f"{source_id}_p{page_num + 1}_c{tc.chunk_index}"
                try:
                    embedding = await embed_text(tc.text)
                    all_vectors.append({
                        "id": chunk_id,
                        "values": embedding,
                        "metadata": {
                            "source_id": source_id,
                            "source_type": "guideline_pdf",
                            "modality": "text",
                            "page_num": page_num + 1,
                            "section_title": tc.section_title or "",
                            "chunk_index": tc.chunk_index,
                            "text": tc.text[:4000],  # Pinecone metadata limit
                            "ingested_at": datetime.now(timezone.utc).isoformat(),
                        },
                    })
                    chunk_count += 1
                except Exception as e:
                    logger.error("Failed to embed text chunk %s: %s", chunk_id, e)

        # --- Image extraction ---
        images = page.get_images(full=True)
        for img_idx, img_info in enumerate(images):
            xref = img_info[0]
            try:
                pix = fitz.Pixmap(doc, xref)
                if pix.n > 4:  # CMYK or other
                    pix = fitz.Pixmap(fitz.csRGB, pix)

                img_bytes = pix.tobytes("png")
                img_name = f"p{page_num + 1}_img{img_idx}.png"
                img_path = _save_image(img_bytes, image_dir, img_name)

                chunk_id = f"{source_id}_p{page_num + 1}_img{img_idx}"
                caption = f"Image from page {page_num + 1} of PDF document"

                try:
                    embedding = await embed_image(img_path, caption=caption)
                    all_vectors.append({
                        "id": chunk_id,
                        "values": embedding,
                        "metadata": {
                            "source_id": source_id,
                            "source_type": "guideline_pdf",
                            "modality": "image",
                            "page_num": page_num + 1,
                            "section_title": "",
                            "chunk_index": img_idx,
                            "image_path": img_path,
                            "caption": caption,
                            "text": caption,
                            "ingested_at": datetime.now(timezone.utc).isoformat(),
                        },
                    })
                    chunk_count += 1
                except Exception as e:
                    logger.error("Failed to embed image %s: %s", chunk_id, e)

                pix = None  # free memory
            except Exception as e:
                logger.error("Failed to extract image xref=%s from page %d: %s", xref, page_num + 1, e)

    doc.close()

    # Upsert to Pinecone in batches of 100
    if all_vectors:
        for i in range(0, len(all_vectors), 100):
            batch = all_vectors[i:i + 100]
            pinecone_index.upsert(vectors=batch)

    logger.info("PDF ingestion complete: %d chunks from %s", chunk_count, file_path)
    return chunk_count


# ---------------------------------------------------------------------------
# PPTX Ingestion
# ---------------------------------------------------------------------------

async def ingest_pptx(
    file_path: str,
    source_id: str,
    pinecone_index,
) -> int:
    """Ingest a PowerPoint file: extract text per slide, embed, store."""
    prs = Presentation(file_path)
    all_vectors = []
    chunk_count = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        # Extract text from all shapes
        slide_texts = []
        slide_title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_texts.append(text)
            if shape.shape_type == 13 and hasattr(shape, "title"):
                # Title placeholder
                pass
        # Try to get slide title
        if slide.shapes.title and slide.shapes.title.text:
            slide_title = slide.shapes.title.text.strip()

        full_text = "\n".join(slide_texts)
        if not full_text.strip():
            continue

        chunk_id = f"{source_id}_slide{slide_num}"
        try:
            embedding = await embed_text(full_text)
            all_vectors.append({
                "id": chunk_id,
                "values": embedding,
                "metadata": {
                    "source_id": source_id,
                    "source_type": "guideline_pptx",
                    "modality": "text",
                    "page_num": slide_num,
                    "section_title": slide_title,
                    "chunk_index": slide_num - 1,
                    "text": full_text[:4000],
                    "ingested_at": datetime.now(timezone.utc).isoformat(),
                },
            })
            chunk_count += 1
        except Exception as e:
            logger.error("Failed to embed slide %d: %s", slide_num, e)

        # Extract images from slide shapes
        image_dir = os.path.join(settings.UPLOAD_DIR, "images", source_id)
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    img_bytes = image.blob
                    ext = image.content_type.split("/")[-1] if image.content_type else "png"
                    img_name = f"slide{slide_num}_img{shape_idx}.{ext}"
                    img_path = _save_image(img_bytes, image_dir, img_name)

                    img_chunk_id = f"{source_id}_slide{slide_num}_img{shape_idx}"
                    caption = f"Image from slide {slide_num}: {slide_title}" if slide_title else f"Image from slide {slide_num}"

                    try:
                        embedding = await embed_image(img_path, caption=caption)
                        all_vectors.append({
                            "id": img_chunk_id,
                            "values": embedding,
                            "metadata": {
                                "source_id": source_id,
                                "source_type": "guideline_pptx",
                                "modality": "image",
                                "page_num": slide_num,
                                "section_title": slide_title,
                                "chunk_index": shape_idx,
                                "image_path": img_path,
                                "caption": caption,
                                "text": caption,
                                "ingested_at": datetime.now(timezone.utc).isoformat(),
                            },
                        })
                        chunk_count += 1
                    except Exception as e:
                        logger.error("Failed to embed slide image %s: %s", img_chunk_id, e)
                except Exception as e:
                    logger.error("Failed to extract image from slide %d: %s", slide_num, e)

    # Upsert to Pinecone
    if all_vectors:
        for i in range(0, len(all_vectors), 100):
            batch = all_vectors[i:i + 100]
            pinecone_index.upsert(vectors=batch)

    logger.info("PPTX ingestion complete: %d chunks from %s", chunk_count, file_path)
    return chunk_count


# ---------------------------------------------------------------------------
# Video Frame Ingestion
# ---------------------------------------------------------------------------

async def ingest_video(
    file_path: str,
    source_id: str,
    pinecone_index,
    fps_sample: float = 1.0,
    similarity_threshold: int = 10,
) -> int:
    """Ingest a video: extract keyframes at fps_sample, deduplicate, embed.

    Uses perceptual hashing to cluster similar frames and keep representatives.
    """
    import cv2
    import numpy as np

    cap = cv2.VideoCapture(file_path)
    if not cap.isOpened():
        raise ValueError(f"Cannot open video: {file_path}")

    video_fps = cap.get(cv2.CAP_PROP_FPS)
    frame_interval = max(1, int(video_fps / fps_sample))

    image_dir = os.path.join(settings.UPLOAD_DIR, "images", source_id)
    _ensure_dir(image_dir)

    # Extract and deduplicate frames
    frames: list[tuple[int, str, np.ndarray]] = []  # (frame_num, path, hash)
    frame_num = 0
    saved_count = 0

    def _phash(img_gray: np.ndarray, hash_size: int = 8) -> np.ndarray:
        """Compute perceptual hash of a grayscale image."""
        resized = cv2.resize(img_gray, (hash_size + 1, hash_size))
        diff = resized[:, 1:] > resized[:, :-1]
        return diff.flatten()

    def _hamming(h1: np.ndarray, h2: np.ndarray) -> int:
        return int(np.sum(h1 != h2))

    prev_hashes: list[np.ndarray] = []

    while True:
        ret, frame = cap.read()
        if not ret:
            break

        if frame_num % frame_interval == 0:
            gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
            h = _phash(gray)

            # Check similarity against previously saved frames
            is_duplicate = False
            for ph in prev_hashes:
                if _hamming(h, ph) < similarity_threshold:
                    is_duplicate = True
                    break

            if not is_duplicate:
                timestamp_sec = frame_num / video_fps
                img_name = f"frame_{frame_num}_{int(timestamp_sec)}s.png"
                img_path = os.path.join(image_dir, img_name)
                cv2.imwrite(img_path, frame)
                frames.append((frame_num, img_path, h))
                prev_hashes.append(h)
                saved_count += 1

        frame_num += 1

    cap.release()

    # Optionally OCR on frames with text
    # (pytesseract is optional — skip if not available)
    ocr_available = False
    try:
        import pytesseract
        ocr_available = True
    except ImportError:
        pass

    # Embed each representative frame
    all_vectors = []
    chunk_count = 0

    for idx, (fnum, img_path, _) in enumerate(frames):
        timestamp_sec = fnum / video_fps
        ocr_text = ""

        if ocr_available:
            try:
                import pytesseract
                img = Image.open(img_path)
                ocr_text = pytesseract.image_to_string(img).strip()
            except Exception:
                pass

        caption = f"Video frame at {int(timestamp_sec)}s"
        if ocr_text:
            caption += f" — Text: {ocr_text[:200]}"

        chunk_id = f"{source_id}_frame{idx}"
        try:
            embedding = await embed_image(img_path, caption=caption)
            all_vectors.append({
                "id": chunk_id,
                "values": embedding,
                "metadata": {
                    "source_id": source_id,
                    "source_type": "guideline_video",
                    "modality": "image",
                    "page_num": 0,
                    "section_title": "",
                    "chunk_index": idx,
                    "image_path": img_path,
                    "caption": caption,
                    "text": ocr_text[:4000] if ocr_text else caption,
                    "timestamp_sec": int(timestamp_sec),
                    "frame_index": fnum,
                    "ingested_at": datetime.now(timezone.utc).isoformat(),
                },
            })
            chunk_count += 1
        except Exception as e:
            logger.error("Failed to embed video frame %s: %s", chunk_id, e)

    # Upsert to Pinecone
    if all_vectors:
        for i in range(0, len(all_vectors), 100):
            batch = all_vectors[i:i + 100]
            pinecone_index.upsert(vectors=batch)

    logger.info("Video ingestion complete: %d frames from %s", chunk_count, file_path)
    return chunk_count


# ---------------------------------------------------------------------------
# Text / Markdown Ingestion
# ---------------------------------------------------------------------------

async def ingest_text(
    file_path: str,
    source_id: str,
    pinecone_index,
) -> int:
    """Ingest a plain text or markdown file."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()

    ext = Path(file_path).suffix.lower()
    if ext in (".md", ".markdown"):
        text_chunks = chunk_markdown(content)
    else:
        text_chunks = chunk_plain_text(content)

    source_type = "guideline_text"
    all_vectors = []
    chunk_count = 0

    for tc in text_chunks:
        chunk_id = f"{source_id}_c{tc.chunk_index}"
        try:
            embedding = await embed_text(tc.text)
            all_vectors.append({
                "id": chunk_id,
                "values": embedding,
                "metadata": {
                    "source_id": source_id,
                    "source_type": source_type,
                    "modality": "text",
                    "page_num": 0,
                    "section_title": tc.section_title or "",
                    "chunk_index": tc.chunk_index,
                    "text": tc.text[:4000],
                    "ingested_at": datetime.now(timezone.utc).isoformat(),
                },
            })
            chunk_count += 1
        except Exception as e:
            logger.error("Failed to embed text chunk %s: %s", chunk_id, e)

    if all_vectors:
        for i in range(0, len(all_vectors), 100):
            batch = all_vectors[i:i + 100]
            pinecone_index.upsert(vectors=batch)

    logger.info("Text ingestion complete: %d chunks from %s", chunk_count, file_path)
    return chunk_count


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS = {
    ".pdf": ("guideline_pdf", ingest_pdf),
    ".pptx": ("guideline_pptx", ingest_pptx),
    ".ppt": ("guideline_pptx", ingest_pptx),
    ".mp4": ("guideline_video", ingest_video),
    ".avi": ("guideline_video", ingest_video),
    ".mov": ("guideline_video", ingest_video),
    ".mkv": ("guideline_video", ingest_video),
    ".webm": ("guideline_video", ingest_video),
    ".txt": ("guideline_text", ingest_text),
    ".md": ("guideline_text", ingest_text),
    ".markdown": ("guideline_text", ingest_text),
}


async def ingest_document(
    file_path: str,
    source_id: str,
    pinecone_index,
) -> tuple[str, int]:
    """Route a document to the appropriate ingestion pipeline.

    Returns (source_type, total_chunks).
    """
    ext = Path(file_path).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {list(SUPPORTED_EXTENSIONS.keys())}")

    source_type, ingest_fn = SUPPORTED_EXTENSIONS[ext]
    total_chunks = await ingest_fn(file_path, source_id, pinecone_index)
    return source_type, total_chunks
